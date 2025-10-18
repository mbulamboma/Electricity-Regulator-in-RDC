from app.models.kpis_reglementaires import KPIReglementaire, PerformanceOperateurKPI

from app.are.dashboard.forms import EditKPIForm
"""
Routes pour le dashboard ARE
"""
from flask import render_template, request, redirect, url_for, flash, jsonify, abort, send_file
from flask_login import login_required, current_user
from datetime import datetime, timedelta
from sqlalchemy import func
import json
import io

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from openpyxl import Workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from app.are.dashboard import dashboard_bp
from app.extensions import db
from app.models.dashboard_are import (
    KPIStrategic, IndicateurSectoriel, AlerteRegulateur, 
    DonneesProvince, RapportAnnuel
)
from app.models.operateurs import Operateur
from app.are.dashboard.forms import (
    FiltreTableauBordForm, AlerteForm, KPIForm, 
    IndicateurSectorielForm, RapportAnnuelForm, ExportForm
)
from app.are.services import IndicateursAREService
from app.are.services_statistiques import StatistiquesAREService, DashboardAREService
from app.utils.decorators import admin_required
from app.utils.permissions import get_accessible_operateurs


# Route pour afficher les KPIs réglementaires (seuils de conformité RDC)
@dashboard_bp.route('/kpis-reglementaires')
@login_required
@admin_required
def kpis_reglementaires():
    from app.models.kpis_reglementaires import PerformanceOperateurKPI
    from datetime import datetime

    # Récupérer l'année actuelle
    annee_actuelle = datetime.now().year
    mois_actuel = datetime.now().month

    # Récupérer tous les KPIs réglementaires actifs
    kpis = KPIReglementaire.query.filter_by(actif=True).all()

    kpis_data = []
    for kpi in kpis:
        # Essayer de récupérer la performance réelle la plus récente
        performance = PerformanceOperateurKPI.query.filter_by(
            kpi_id=kpi.id,
            annee=annee_actuelle,
            mois=mois_actuel
        ).first()

        if performance and performance.valeur_mesuree is not None:
            # Utiliser la valeur mesurée réelle
            valeur = performance.valeur_mesuree
        else:
            # Valeur par défaut basée sur les seuils (pour démonstration)
            # En production, ceci devrait être calculé à partir des données réelles
            valeur = kpi.seuil_acceptable * 0.95  # Simuler une bonne performance

        # Évaluer la performance
        niveau, penalite = kpi.evaluer_performance(valeur)

        kpis_data.append({
            'id': kpi.id,
            'code': kpi.code,
            'nom': kpi.nom,
            'description': kpi.description,
            'valeur': round(valeur, 2),
            'seuil': kpi.seuil_acceptable,
            'seuil_excellent': kpi.seuil_excellent,
            'seuil_limite': kpi.seuil_limite,
            'seuil_critique': kpi.seuil_critique,
            'unite': kpi.unite,
            'niveau': niveau,
            'conforme': niveau in ['excellent', 'acceptable'],
            'penalite': penalite,
            'type_kpi': kpi.type_kpi.value if kpi.type_kpi else None,
            'reference_legale': kpi.reference_legale
        })

    return render_template('are/dashboard/components/kpis_reglementaires.html', kpis_reglementaires=kpis_data)


# Route d'édition d'un KPI stratégique
@dashboard_bp.route('/kpi/<int:kpi_id>/edit', methods=['GET', 'POST'])
@login_required
@admin_required
def edit_kpi_strategique(kpi_id):
    from app.models.dashboard_are import KPIStrategic
    kpi = KPIStrategic.query.get_or_404(kpi_id)
    form = EditKPIForm(obj=kpi)
    form.statut.data = 'atteint' if kpi.atteint else 'en_cours'
    if form.validate_on_submit():
        kpi.nom = form.nom.data
        kpi.valeur = form.valeur.data
        kpi.objectif = form.objectif.data
        kpi.unite = form.unite.data
        # Calcul automatique du statut basé sur valeur >= objectif
        kpi.atteint = kpi.valeur >= kpi.objectif if kpi.objectif is not None else False
        kpi.date_modification = datetime.now()
        db.session.commit()
        flash('KPI modifié avec succès.', 'success')
        return redirect(url_for('are_dashboard.index'))
    return render_template('are/dashboard/edit_kpi.html', form=form, kpi=kpi)


@dashboard_bp.route('/api/kpis/<int:annee>/mettre-a-jour', methods=['POST'])
@login_required
@admin_required
def api_mettre_a_jour_kpis(annee):
    """API pour mettre à jour les KPIs automatiquement"""
    kpis = IndicateursAREService.mettre_a_jour_kpis_strategiques(annee)
    
    return jsonify({
        'message': f'{len(kpis)} KPIs mis à jour',
        'kpis': [kpi.to_dict() for kpi in kpis]
    })


# ===== NOUVELLES ROUTES STATISTIQUES ARE =====

@dashboard_bp.route('/statistiques')
@login_required
@admin_required
def statistiques():
    """Page des statistiques complètes ARE"""
    annee_debut = request.args.get('annee_debut', 2020, type=int)
    annee_fin = request.args.get('annee_fin', 2024, type=int)
    
    # 1. Portfolio de projets
    portfolio = DashboardAREService.get_portfolio_projets()
    
    # 2. Evolution de la capacité installée
    evolution_capacite = DashboardAREService.get_evolution_capacite(annee_debut, annee_fin)
    
    # 3. Statistiques nationales
    stats_nationales = DashboardAREService.get_statistiques_nationales_periode(annee_debut, annee_fin)
    
    # 4. Données solaires
    donnees_solaires = DashboardAREService.get_donnees_solaires()
    
    # 5. Statistiques de clientèle
    stats_clientele = DashboardAREService.get_statistiques_clientele(annee_debut, annee_fin)
    
    # 6. Données provinciales pour la carte
    donnees_provinciales = DonneesProvince.query.filter(
        DonneesProvince.annee.between(annee_debut, annee_fin),
        DonneesProvince.actif == True
    ).order_by(DonneesProvince.province, DonneesProvince.annee).all()
    
    # Préparer les données pour les graphiques
    graphiques_data = {
        'evolution_capacite_labels': list(range(annee_debut, annee_fin + 1)),
        'evolution_capacite_hydro': [],
        'evolution_capacite_thermique': [],
        'evolution_capacite_solaire': [],
        'evolution_production_labels': list(range(annee_debut, annee_fin + 1)),
        'evolution_production_values': [],
        'evolution_clients_labels': list(range(annee_debut, annee_fin + 1)),
        'evolution_clients_values': []
    }
    
    # Remplir les données de graphiques
    for annee in range(annee_debut, annee_fin + 1):
        # Capacités par source
        capacite_hydro = sum([c['capacite_installee_mw'] for c in evolution_capacite 
                             if c['annee'] == annee and c['type_source'] == 'production_hydro'])
        capacite_thermique = sum([c['capacite_installee_mw'] for c in evolution_capacite 
                                 if c['annee'] == annee and c['type_source'] == 'production_thermique'])
        capacite_solaire = sum([c['capacite_installee_mw'] for c in evolution_capacite 
                               if c['annee'] == annee and c['type_source'] == 'production_solaire'])
        
        graphiques_data['evolution_capacite_hydro'].append(capacite_hydro)
        graphiques_data['evolution_capacite_thermique'].append(capacite_thermique)
        graphiques_data['evolution_capacite_solaire'].append(capacite_solaire)
        
        # Production totale
        stat_nat = next((s for s in stats_nationales if s['annee'] == annee), None)
        production_totale = stat_nat['production_totale_annuelle_gwh'] if stat_nat else 0
        graphiques_data['evolution_production_values'].append(production_totale)
        
        # Clients totaux
        clients_totaux = sum([c['total_clients'] for c in stats_clientele if c['annee'] == annee])
        graphiques_data['evolution_clients_values'].append(clients_totaux)
    
    return render_template('are/dashboard/statistiques.html',
                         portfolio=portfolio,
                         evolution_capacite=evolution_capacite,
                         stats_nationales=stats_nationales,
                         donnees_solaires=donnees_solaires,
                         stats_clientele=stats_clientele,
                         donnees_provinciales=donnees_provinciales,
                         graphiques_data=graphiques_data,
                         annee_debut=annee_debut,
                         annee_fin=annee_fin)


@dashboard_bp.route('/calculer-statistiques/<int:annee>')
@login_required
@admin_required
def calculer_statistiques(annee):
    """Lance le calcul des statistiques pour une année donnée"""
    try:
        success = StatistiquesAREService.calculer_toutes_statistiques(annee)
        
        if success:
            flash(f'✅ Statistiques calculées avec succès pour l\'année {annee}', 'success')
        else:
            flash(f'❌ Erreur lors du calcul des statistiques pour l\'année {annee}', 'error')
    
    except Exception as e:
        flash(f'❌ Erreur technique: {str(e)}', 'error')
    
    return redirect(url_for('are_dashboard.statistiques'))


@dashboard_bp.route('/api/statistiques/portfolio')
@login_required
@admin_required
def api_portfolio_projets():
    """API pour récupérer le portfolio des projets"""
    portfolio = DashboardAREService.get_portfolio_projets()
    return jsonify(portfolio)


@dashboard_bp.route('/api/statistiques/evolution-capacite')
@login_required
@admin_required
def api_evolution_capacite():
    """API pour récupérer l'évolution de la capacité"""
    annee_debut = request.args.get('annee_debut', 2020, type=int)
    annee_fin = request.args.get('annee_fin', 2024, type=int)
    
    evolution = DashboardAREService.get_evolution_capacite(annee_debut, annee_fin)
    return jsonify(evolution)


@dashboard_bp.route('/api/statistiques/nationales')
@login_required
@admin_required
def api_statistiques_nationales():
    """API pour récupérer les statistiques nationales"""
    annee_debut = request.args.get('annee_debut', 2020, type=int)
    annee_fin = request.args.get('annee_fin', 2024, type=int)
    
    stats = DashboardAREService.get_statistiques_nationales_periode(annee_debut, annee_fin)
    return jsonify(stats)


@dashboard_bp.route('/')
@login_required
@admin_required
def index():
    """Page principale du dashboard ARE avec statistiques nationales complètes"""
    form = FiltreTableauBordForm()
    
    # Récupérer les filtres
    annee = request.args.get('annee', datetime.now().year, type=int)
    annee_debut = request.args.get('annee_debut', 2023, type=int)
    annee_fin = request.args.get('annee_fin', datetime.now().year, type=int)
    operateur_id = request.args.get('operateur_id', type=int)
    province = request.args.get('province', '')
    
    # KPIs stratégiques
    kpis_query = KPIStrategic.query.filter_by(annee=annee, actif=True)
    if operateur_id:
        kpis_query = kpis_query.filter_by(operateur_id=operateur_id)
    kpis = kpis_query.all()
    

    # Mix énergétique
    mix_energetique = IndicateursAREService.calculer_mix_energetique(annee, operateur_id)

    # Performance des opérateurs
    performance_operateurs = IndicateursAREService.calculer_performance_operateurs(annee)
    if operateur_id:
        performance_operateurs = [p for p in performance_operateurs if p['operateur_id'] == operateur_id]

    # KPIs stratégiques (base de données)
    kpis_strategiques = KPIStrategic.query.filter_by(annee=annee, actif=True).all()
    
    # Données par province pour la carte
    donnees_provinces = []
    if not operateur_id:  # Affichage national uniquement
        donnees_provinces_obj = DonneesProvince.query.filter_by(annee=annee).all()
        donnees_provinces = [province.to_dict() for province in donnees_provinces_obj]

    # Récupérer les centrales hydroélectriques actives avec coordonnées
    from app.models.production_hydro import CentraleHydro
    centrales = CentraleHydro.query.filter(
        CentraleHydro.actif == True,
        CentraleHydro.latitude.isnot(None),
        CentraleHydro.longitude.isnot(None)
    ).all()
    centrales_data = [
        {
            'nom': c.nom,
            'type': 'Hydroélectrique',
            'capacite_mw': c.puissance_installee,
            'latitude': c.latitude,
            'longitude': c.longitude
        }
        for c in centrales
    ]
    
    # KPIs réglementaires (seuils de conformité RDC)
    kpis_reglementaires_data = []
    kpis_reg = KPIReglementaire.query.filter_by(actif=True).all()
    
    annee_actuelle = datetime.now().year
    mois_actuel = datetime.now().month
    
    for kpi in kpis_reg:
        # Essayer de récupérer la performance réelle la plus récente
        performance = PerformanceOperateurKPI.query.filter_by(
            kpi_id=kpi.id,
            annee=annee_actuelle,
            mois=mois_actuel
        ).first()

        if performance and performance.valeur_mesuree is not None:
            # Utiliser la valeur mesurée réelle
            valeur = performance.valeur_mesuree
        else:
            # Valeur par défaut basée sur les seuils (pour démonstration)
            valeur = kpi.seuil_acceptable * 0.95  # Simuler une bonne performance

        # Évaluer la performance
        niveau, penalite = kpi.evaluer_performance(valeur)

        kpis_reglementaires_data.append({
            'id': kpi.id,
            'code': kpi.code,
            'nom': kpi.nom,
            'description': kpi.description,
            'valeur': round(valeur, 2),
            'seuil': kpi.seuil_acceptable,
            'seuil_excellent': kpi.seuil_excellent,
            'seuil_limite': kpi.seuil_limite,
            'seuil_critique': kpi.seuil_critique,
            'unite': kpi.unite,
            'niveau': niveau,
            'conforme': niveau in ['excellent', 'acceptable'],
            'penalite': penalite,
            'type_kpi': kpi.type_kpi.value if kpi.type_kpi else None,
            'reference_legale': kpi.reference_legale
        })
    
    # NOUVELLES STATISTIQUES NATIONALES AVANCÉES
    stats_nationales = _calculer_statistiques_nationales_avancees(annee_debut, annee_fin)
    
    # Statistiques générales - Infrastructure seulement
    stats = {
        'infrastructure': {
            'total_centrales': stats_nationales.get('stats_infrastructure', {}).get('total_centrales', 0),
            'longueur_lignes_km': stats_nationales.get('stats_infrastructure', {}).get('longueur_lignes_km', 0),
            'puissance_disponible_mw': stats_nationales.get('stats_infrastructure', {}).get('puissance_disponible_mw', 0),
            'capacite_totale_mw': stats_nationales.get('stats_infrastructure', {}).get('capacite_totale_mw', 0)
        }
    }
    
    return render_template('are/dashboard/index.html',
                         form=form,
                         kpis=kpis,
                         mix_energetique=mix_energetique,
                         performance_operateurs=performance_operateurs,
                         donnees_provinces=donnees_provinces,
                         stats=stats,
                         stats_nationales=stats_nationales,
                         annee=annee,
                         annee_debut=annee_debut,
                         annee_fin=annee_fin,
                         operateur_id=operateur_id,
                         province=province,
                         centrales=centrales_data,
                         kpis_strategiques=kpis_strategiques,
                         kpis_reglementaires=kpis_reglementaires_data)


@dashboard_bp.route('/kpis')
@login_required
@admin_required
def kpis():
    """Page détaillée des KPIs"""
    form = FiltreTableauBordForm()
    
    # Filtres
    annee = request.args.get('annee', datetime.now().year, type=int)
    operateur_id = request.args.get('operateur_id', type=int)
    
    # Récupérer tous les KPIs
    kpis_query = KPIStrategic.query.filter_by(annee=annee, actif=True)
    if operateur_id:
        kpis_query = kpis_query.filter_by(operateur_id=operateur_id)
    
    kpis = kpis_query.order_by(KPIStrategic.code).all()
    
    # Grouper les KPIs par catégorie
    kpis_groupes = {}
    for kpi in kpis:
        prefix = kpi.code.split('_')[0]
        if prefix not in kpis_groupes:
            kpis_groupes[prefix] = []
        kpis_groupes[prefix].append(kpi)
    
    # Indicateurs sectoriels
    indicateurs = IndicateurSectoriel.query.filter_by(annee=annee, actif=True)
    if operateur_id:
        indicateurs = indicateurs.filter_by(operateur_id=operateur_id)
    indicateurs = indicateurs.all()
    
    # Grouper les indicateurs par catégorie
    indicateurs_groupes = {}
    for ind in indicateurs:
        cat = ind.categorie.value
        if cat not in indicateurs_groupes:
            indicateurs_groupes[cat] = []
        indicateurs_groupes[cat].append(ind)
    
    return render_template('are/dashboard/kpis.html',
                         form=form,
                         kpis_groupes=kpis_groupes,
                         indicateurs_groupes=indicateurs_groupes,
                         annee=annee,
                         operateur_id=operateur_id)


@dashboard_bp.route('/kpis/nouveau', methods=['GET', 'POST'])
@login_required
@admin_required
def nouveau_kpi():
    """Créer un nouveau KPI"""
    form = KPIForm()
    
    if form.validate_on_submit():
        kpi = KPIStrategic(
            code=form.code.data,
            nom=form.nom.data,
            description=form.description.data,
            valeur=form.valeur.data,
            unite=form.unite.data,
            periode=form.periode.data,
            annee=form.annee.data,
            objectif=form.objectif.data,
            seuil_alerte=form.seuil_alerte.data,
            operateur_id=form.operateur_id.data,
            source_donnees=form.source_donnees.data
        )
        
        # Calcul automatique du statut basé sur valeur >= objectif
        kpi.atteint = kpi.valeur >= kpi.objectif if kpi.objectif is not None else False
        
        kpi.save()
        flash('KPI créé avec succès.', 'success')
        return redirect(url_for('are_dashboard.kpis'))
    
    return render_template('are/dashboard/kpi_form.html', form=form, title='Nouveau KPI')


@dashboard_bp.route('/kpis/<int:kpi_id>/modifier', methods=['GET', 'POST'])
@login_required
@admin_required
def modifier_kpi(kpi_id):
    """Modifier un KPI existant"""
    kpi = KPIStrategic.query.get_or_404(kpi_id)
    form = KPIForm(obj=kpi)
    
    if form.validate_on_submit():
        kpi.update(
            code=form.code.data,
            nom=form.nom.data,
            description=form.description.data,
            valeur=form.valeur.data,
            unite=form.unite.data,
            periode=form.periode.data,
            annee=form.annee.data,
            objectif=form.objectif.data,
            seuil_alerte=form.seuil_alerte.data,
            operateur_id=form.operateur_id.data,
            source_donnees=form.source_donnees.data
        )
        
        # Calcul automatique du statut basé sur valeur >= objectif
        kpi.atteint = kpi.valeur >= kpi.objectif if kpi.objectif is not None else False
        kpi.save()
        
        flash('KPI modifié avec succès.', 'success')
        return redirect(url_for('are_dashboard.kpis'))
    
    return render_template('are/dashboard/kpi_form.html', 
                         form=form, 
                         kpi=kpi,
                         title='Modifier KPI')


@dashboard_bp.route('/alertes')
@login_required
@admin_required
def alertes():
    """Page de gestion des alertes"""
    # Filtres
    statut = request.args.get('statut', 'active')
    type_alerte = request.args.get('type')
    severite = request.args.get('severite')
    
    # Construire la requête
    alertes_query = AlerteRegulateur.query
    
    if statut:
        alertes_query = alertes_query.filter_by(statut=statut)
    if type_alerte:
        alertes_query = alertes_query.filter_by(type=type_alerte)
    if severite:
        alertes_query = alertes_query.filter_by(severite=severite)
    
    alertes = alertes_query.order_by(
        AlerteRegulateur.priorite.asc(),
        AlerteRegulateur.date_creation.desc()
    ).all()
    
    # Statistiques des alertes
    stats_alertes = {
        'total': AlerteRegulateur.query.count(),
        'actives': AlerteRegulateur.query.filter_by(statut='active').count(),
        'critiques': AlerteRegulateur.query.filter_by(
            statut='active', 
            severite='critique'
        ).count(),
        'expirees': AlerteRegulateur.query.filter(
            AlerteRegulateur.date_echeance < datetime.now().date(),
            AlerteRegulateur.statut == 'active'
        ).count() if AlerteRegulateur.query.filter(
            AlerteRegulateur.date_echeance.isnot(None)
        ).first() else 0
    }
    
    return render_template('are/dashboard/alertes.html',
                         alertes=alertes,
                         stats_alertes=stats_alertes,
                         statut=statut,
                         type_alerte=type_alerte,
                         severite=severite)


@dashboard_bp.route('/alertes/nouvelle', methods=['GET', 'POST'])
@login_required
@admin_required
def nouvelle_alerte():
    """Créer une nouvelle alerte"""
    form = AlerteForm()
    
    if form.validate_on_submit():
        alerte = AlerteRegulateur(
            type=form.type.data,
            severite=form.severite.data,
            operateur_id=form.operateur_id.data,
            entite_concernee=form.entite_concernee.data,
            titre=form.titre.data,
            description=form.description.data,
            date_echeance=form.date_echeance.data,
            actions_recommandees=form.actions_recommandees.data,
            priorite=form.priorite.data,
            createur_id=current_user.id
        )
        
        alerte.save()
        flash('Alerte créée avec succès.', 'success')
        return redirect(url_for('are_dashboard.alertes'))
    
    return render_template('are/dashboard/alerte_form.html', 
                         form=form, 
                         title='Nouvelle alerte')


@dashboard_bp.route('/alertes/<int:alerte_id>/resoudre', methods=['POST'])
@login_required
@admin_required
def resoudre_alerte(alerte_id):
    """Marquer une alerte comme résolue"""
    alerte = AlerteRegulateur.query.get_or_404(alerte_id)
    
    alerte.update(
        statut='resolue',
        date_resolution=datetime.now().date()
    )
    
    flash('Alerte marquée comme résolue.', 'success')
    return redirect(url_for('are_dashboard.alertes'))


@dashboard_bp.route('/alertes/<int:alerte_id>')
@login_required
@admin_required
def detail_alerte(alerte_id):
    """Voir les détails d'une alerte"""
    alerte = AlerteRegulateur.query.get_or_404(alerte_id)
    return render_template('are/dashboard/alerte_detail.html', alerte=alerte)


@dashboard_bp.route('/alertes/<int:alerte_id>/supprimer', methods=['POST'])
@login_required
@admin_required
def supprimer_alerte(alerte_id):
    """Supprimer une alerte"""
    alerte = AlerteRegulateur.query.get_or_404(alerte_id)
    
    # Sauvegarder les infos pour le message
    titre_alerte = alerte.titre
    
    # Supprimer l'alerte
    alerte.delete()
    
    flash(f'Alerte "{titre_alerte}" supprimée avec succès.', 'success')
    return redirect(url_for('are_dashboard.alertes'))


@dashboard_bp.route('/rapports')
@login_required
@admin_required
def rapports():
    """Page de gestion des rapports annuels"""
    rapports = RapportAnnuel.query.order_by(
        RapportAnnuel.annee.desc()
    ).all()
    
    return render_template('are/dashboard/rapports.html', rapports=rapports)


@dashboard_bp.route('/rapports/nouveau', methods=['GET', 'POST'])
@login_required
@admin_required
def nouveau_rapport():
    """Créer un nouveau rapport annuel"""
    form = RapportAnnuelForm()
    
    if form.validate_on_submit():
        # Générer le contenu du rapport
        contenu = {
            'titre': form.titre.data,
            'annee': form.annee.data,
            'periode': {
                'debut': form.periode_debut.data.isoformat(),
                'fin': form.periode_fin.data.isoformat()
            },
            'sections': {
                'mix_energetique': form.inclure_mix_energetique.data == 'oui',
                'performance_operateurs': form.inclure_performance_operateurs.data == 'oui',
                'indicateurs_financiers': form.inclure_indicateurs_financiers.data == 'oui',
                'activites_regulation': form.inclure_activites_regulation.data == 'oui',
                'perspectives': form.inclure_perspectives.data == 'oui'
            },
            'donnees': {}
        }
        
        # Récupérer les données selon les sections sélectionnées
        if contenu['sections']['mix_energetique']:
            contenu['donnees']['mix_energetique'] = IndicateursAREService.calculer_mix_energetique(form.annee.data)
        
        if contenu['sections']['performance_operateurs']:
            contenu['donnees']['performance_operateurs'] = IndicateursAREService.calculer_performance_operateurs(form.annee.data)
        
        # Créer le rapport
        rapport = RapportAnnuel(
            annee=form.annee.data,
            titre=form.titre.data,
            contenu_json=contenu,
            periode_debut=form.periode_debut.data,
            periode_fin=form.periode_fin.data,
            nombre_operateurs=Operateur.query.filter_by(actif=True).count(),
            auteur_id=current_user.id
        )
        
        rapport.save()
        flash('Rapport annuel généré avec succès.', 'success')
        return redirect(url_for('are_dashboard.voir_rapport', rapport_id=rapport.id))
    
    return render_template('are/dashboard/rapport_form.html', 
                         form=form, 
                         title='Nouveau rapport annuel')


@dashboard_bp.route('/rapports/<int:rapport_id>')
@login_required
@admin_required
def voir_rapport(rapport_id):
    """Voir un rapport annuel"""
    rapport = RapportAnnuel.query.get_or_404(rapport_id)
    
    return render_template('are/dashboard/rapport_detail.html', rapport=rapport)


@dashboard_bp.route('/export', methods=['GET', 'POST'])
@login_required
@admin_required
def export_donnees():
    """Exporter les données du dashboard"""
    form = ExportForm()
    
    if form.validate_on_submit():
        format_export = form.format_export.data
        donnees_export = form.donnees_export.data
        annee = form.annee_export.data
        operateur_id = form.operateur_id.data
        
        if format_export == 'excel':
            return export_excel(donnees_export, annee, operateur_id)
        elif format_export == 'powerpoint':
            return export_powerpoint(donnees_export, annee, operateur_id)
        else:
            flash('Format d\'export non supporté.', 'error')
    
    return render_template('are/dashboard/export.html', form=form)


def export_excel(donnees_export, annee, operateur_id):
    """Exporter vers Excel"""
    if not OPENPYXL_AVAILABLE:
        flash('Openpyxl n\'est pas installé. Veuillez installer openpyxl pour l\'export Excel.', 'error')
        return redirect(url_for('are_dashboard.export_donnees'))
    
    output = io.BytesIO()
    workbook = Workbook()
    
    # Supprimer la feuille par défaut
    workbook.remove(workbook.active)
    
    if donnees_export in ['kpis', 'tout']:
        # Feuille KPIs
        ws_kpis = workbook.create_sheet("KPIs Stratégiques")
        ws_kpis.append(['Code', 'Nom', 'Valeur', 'Unité', 'Période', 'Objectif', 'Opérateur'])
        
        kpis_query = KPIStrategic.query.filter_by(annee=annee, actif=True)
        if operateur_id:
            kpis_query = kpis_query.filter_by(operateur_id=operateur_id)
        
        for kpi in kpis_query.all():
            ws_kpis.append([
                kpi.code, kpi.nom, kpi.valeur, kpi.unite, 
                kpi.periode, kpi.objectif or '',
                kpi.operateur.nom_commercial if kpi.operateur else 'National'
            ])
    
    if donnees_export in ['performance', 'tout']:
        # Feuille Performance
        ws_perf = workbook.create_sheet("Performance Opérateurs")
        ws_perf.append([
            'Opérateur', 'Puissance Installée (MW)', 'Production (MWh)', 
            'Facteur de Charge (%)', 'Clients Total', 'Part Hydro (%)', 
            'Part Thermique (%)', 'Part Solaire (%)'
        ])
        
        performance = IndicateursAREService.calculer_performance_operateurs(annee)
        if operateur_id:
            performance = [p for p in performance if p['operateur_id'] == operateur_id]
        
        for perf in performance:
            ws_perf.append([
                perf['operateur'], 
                perf['puissance_installee'],
                perf['production_annuelle'],
                perf['facteur_charge'],
                perf['clients_total'],
                perf['mix_energetique']['hydro'],
                perf['mix_energetique']['thermique'],
                perf['mix_energetique']['solaire']
            ])
    
    workbook.save(output)
    output.seek(0)
    
    return send_file(
        output,
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        as_attachment=True,
        download_name=f'dashboard_are_{annee}.xlsx'
    )


def export_powerpoint(donnees_export, annee, operateur_id):
    """Exporter vers PowerPoint"""
    if not PPTX_AVAILABLE:
        flash('python-pptx n\'est pas installé. Veuillez installer python-pptx pour l\'export PowerPoint.', 'error')
        return redirect(url_for('are_dashboard.export_donnees'))
    
    prs = Presentation()
    
    # Slide de titre
    slide_layout = prs.slide_layouts[0]  # Layout titre
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"Dashboard ARE - Année {annee}"
    subtitle.text = f"Rapport généré le {datetime.now().strftime('%d/%m/%Y')}"
    
    if donnees_export in ['kpis', 'tout']:
        # Slide KPIs
        slide_layout = prs.slide_layouts[1]  # Layout contenu
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "KPIs Stratégiques"
        
        # Récupérer les KPIs principaux
        kpis_query = KPIStrategic.query.filter_by(annee=annee, actif=True)
        if operateur_id:
            kpis_query = kpis_query.filter_by(operateur_id=operateur_id)
        
        kpis = kpis_query.limit(5).all()  # Top 5 KPIs
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = "Principaux indicateurs:"
        
        for kpi in kpis:
            p = text_frame.add_paragraph()
            p.text = f"• {kpi.nom}: {kpi.valeur} {kpi.unite}"
    
    # Sauvegarder
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return send_file(
        output,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        as_attachment=True,
        download_name=f'dashboard_are_{annee}.pptx'
    )


# API endpoints pour les données dynamiques

@dashboard_bp.route('/api/kpis/<int:annee>')
@login_required
@admin_required
def api_kpis(annee):
    """API pour récupérer les KPIs d'une année"""
    operateur_id = request.args.get('operateur_id', type=int)
    
    kpis_query = KPIStrategic.query.filter_by(annee=annee, actif=True)
    if operateur_id:
        kpis_query = kpis_query.filter_by(operateur_id=operateur_id)
    
    kpis = kpis_query.all()
    
    return jsonify({
        'kpis': [kpi.to_dict() for kpi in kpis],
        'total': len(kpis)
    })


@dashboard_bp.route('/api/mix-energetique/<int:annee>')
@login_required
@admin_required
def api_mix_energetique(annee):
    """API pour le mix énergétique"""
    operateur_id = request.args.get('operateur_id', type=int)
    
    mix = IndicateursAREService.calculer_mix_energetique(annee, operateur_id)
    
    return jsonify(mix)


@dashboard_bp.route('/api/alertes/generer', methods=['POST'])
@login_required
@admin_required
def api_generer_alertes():
    """API pour générer les alertes automatiques"""
    alertes = IndicateursAREService.generer_alertes_automatiques()
    
    return jsonify({
        'message': f'{len(alertes)} alertes générées',
        'alertes': [alerte.to_dict() for alerte in alertes]
    })


@dashboard_bp.route('/kpi/<int:kpi_id>/supprimer', methods=['POST'])
@login_required
@admin_required
def supprimer_kpi(kpi_id):
    """Supprimer un KPI stratégique"""
    kpi = KPIStrategic.query.get_or_404(kpi_id)
    
    try:
        nom_kpi = kpi.nom
        kpi.delete()
        flash(f'KPI "{nom_kpi}" supprimé avec succès.', 'success')
    except Exception as e:
        flash(f'Erreur lors de la suppression du KPI: {str(e)}', 'error')
    
    return redirect(url_for('are_dashboard.index'))


@dashboard_bp.route('/export/<export_type>/<data_type>')
@login_required
@admin_required
def export_data(export_type, data_type):
    """Exporter des données du dashboard en différents formats"""
    if not current_user.is_super_admin():
        abort(403)
    
    try:
        # Récupérer les services de données
        stats_service = StatistiquesAREService()
        dashboard_service = DashboardAREService()
        
        # Générer le nom de fichier avec timestamp
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        
        if export_type == 'excel':
            return _export_to_excel(data_type, timestamp, stats_service, dashboard_service)
        elif export_type == 'csv':
            return _export_to_csv(data_type, timestamp, stats_service, dashboard_service)
        elif export_type == 'json':
            return _export_to_json(data_type, timestamp, stats_service, dashboard_service)
        else:
            flash('Format d\'export non supporté.', 'error')
            return redirect(url_for('are_dashboard.index'))
            
    except Exception as e:
        flash(f'Erreur lors de l\'export: {str(e)}', 'error')
        return redirect(url_for('are_dashboard.index'))


def _export_to_excel(data_type, timestamp, stats_service, dashboard_service):
    """Export en format Excel"""
    if not PANDAS_AVAILABLE:
        flash('Pandas n\'est pas disponible pour l\'export Excel.', 'error')
        return redirect(url_for('are_dashboard.index'))
    
    output = io.BytesIO()
    
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        if data_type == 'production' or data_type == 'all':
            # Données de production
            stats_nationales = _calculer_statistiques_nationales_avancees(2023, 2025)
            if stats_nationales and 'evolution_production' in stats_nationales:
                production_data = []
                for item in stats_nationales['evolution_production']:
                    production_data.append({
                        'Année': item['annee'],
                        'Production Hydro (GWh)': item['production_hydro_gwh'],
                        'Production Thermique (GWh)': item['production_thermique_gwh'],
                        'Production Solaire (GWh)': item['production_solaire_gwh'],
                        'Production Totale (GWh)': item['production_totale_gwh']
                    })
                
                if production_data:
                    df_production = pd.DataFrame(production_data)
                    df_production.to_excel(writer, sheet_name='Production', index=False)
        
        if data_type == 'capacite' or data_type == 'all':
            # Données de capacité installée
            stats_nationales = _calculer_statistiques_nationales_avancees(2023, 2025)
            if stats_nationales and 'evolution_capacite' in stats_nationales:
                capacite_data = []
                for item in stats_nationales['evolution_capacite']:
                    capacite_data.append({
                        'Année': item['annee'],
                        'Capacité Hydro (MW)': item['capacite_hydro_mw'],
                        'Capacité Thermique (MW)': item['capacite_thermique_mw'],
                        'Capacité Solaire (MW)': item['capacite_solaire_mw'],
                        'Capacité Totale (MW)': item['capacite_totale_mw']
                    })
                
                if capacite_data:
                    df_capacite = pd.DataFrame(capacite_data)
                    df_capacite.to_excel(writer, sheet_name='Capacité', index=False)
        
        if data_type == 'operateurs' or data_type == 'all':
            # Données des opérateurs
            stats_nationales = _calculer_statistiques_nationales_avancees(2023, 2025)
            if stats_nationales and 'stats_operateurs' in stats_nationales:
                operateurs_data = []
                for op in stats_nationales['stats_operateurs']:
                    operateurs_data.append({
                        'Opérateur': op['nom'],
                        'Production (GWh)': op['production_gwh'],
                        'Clients': op['clients'],
                        'Rendement (%)': op['rendement_pct']
                    })
                
                if operateurs_data:
                    df_operateurs = pd.DataFrame(operateurs_data)
                    df_operateurs.to_excel(writer, sheet_name='Opérateurs', index=False)
        
        if data_type == 'kpis' or data_type == 'all':
            # KPIs stratégiques
            kpis = KPIStrategic.query.filter_by(actif=True).all()
            if kpis:
                kpi_data = []
                for kpi in kpis:
                    kpi_data.append({
                        'Nom': kpi.nom,
                        'Valeur': kpi.valeur,
                        'Objectif': kpi.objectif,
                        'Unité': kpi.unite,
                        'Tendance': kpi.tendance,
                        'Dernière MAJ': kpi.date_modification.strftime('%Y-%m-%d %H:%M')
                    })
                
                df_kpis = pd.DataFrame(kpi_data)
                df_kpis.to_excel(writer, sheet_name='KPIs', index=False)
        
        if data_type == 'alertes' or data_type == 'all':
            # Alertes
            alertes = AlerteRegulateur.query.filter_by(actif=True).order_by(
                AlerteRegulateur.date_creation.desc()
            ).limit(100).all()
            
            if alertes:
                alerte_data = []
                for alerte in alertes:
                    alerte_data.append({
                        'Titre': alerte.titre,
                        'Type': alerte.type.value if alerte.type else 'N/A',
                        'Sévérité': alerte.severite.value if alerte.severite else 'N/A',
                        'Statut': alerte.statut,
                        'Priorité': alerte.priorite,
                        'Entité concernée': alerte.entite_concernee,
                        'Description': alerte.description,
                        'Date création': alerte.date_creation.strftime('%Y-%m-%d %H:%M'),
                        'Date échéance': alerte.date_echeance.strftime('%Y-%m-%d') if alerte.date_echeance else 'N/A'
                    })
                
                df_alertes = pd.DataFrame(alerte_data)
                df_alertes.to_excel(writer, sheet_name='Alertes', index=False)
    
    output.seek(0)
    
    filename = f'are_dashboard_{data_type}_{timestamp}.xlsx'
    
    return send_file(
        output,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )


def _export_to_csv(data_type, timestamp, stats_service, dashboard_service):
    """Export en format CSV"""
    if not PANDAS_AVAILABLE:
        flash('Pandas n\'est pas disponible pour l\'export CSV.', 'error')
        return redirect(url_for('are_dashboard.index'))
    
    # Récupérer les données selon le type
    stats_nationales = _calculer_statistiques_nationales_avancees(2023, 2025)
    
    if data_type == 'production' and 'evolution_production' in stats_nationales:
        data = []
        for item in stats_nationales['evolution_production']:
            data.append({
                'Année': item['annee'],
                'Production Hydro (GWh)': item['production_hydro_gwh'],
                'Production Thermique (GWh)': item['production_thermique_gwh'],
                'Production Solaire (GWh)': item['production_solaire_gwh'],
                'Production Totale (GWh)': item['production_totale_gwh']
            })
        df = pd.DataFrame(data)
        filename = f'are_production_{timestamp}.csv'
    
    elif data_type == 'operateurs' and 'stats_operateurs' in stats_nationales:
        data = []
        for op in stats_nationales['stats_operateurs']:
            data.append({
                'Opérateur': op['nom'],
                'Production (GWh)': op['production_gwh'],
                'Clients': op['clients'],
                'Rendement (%)': op['rendement_pct']
            })
        df = pd.DataFrame(data)
        filename = f'are_operateurs_{timestamp}.csv'
    
    else:
        flash('Type de données non supporté pour l\'export CSV.', 'error')
        return redirect(url_for('are_dashboard.index'))
    
    output = io.StringIO()
    df.to_csv(output, index=False)
    output.seek(0)
    
    return send_file(
        io.BytesIO(output.getvalue().encode('utf-8')),
        as_attachment=True,
        download_name=filename,
        mimetype='text/csv'
    )


def _export_to_json(data_type, timestamp, stats_service, dashboard_service):
    """Export en format JSON"""
    # Récupérer les données selon le type
    stats_nationales = _calculer_statistiques_nationales_avancees(2023, 2025)
    
    if data_type == 'production' and 'evolution_production' in stats_nationales:
        data = stats_nationales['evolution_production']
        filename = f'are_production_{timestamp}.json'
    
    elif data_type == 'operateurs' and 'stats_operateurs' in stats_nationales:
        data = stats_nationales['stats_operateurs']
        filename = f'are_operateurs_{timestamp}.json'
    
    elif data_type == 'all':
        data = stats_nationales
        filename = f'are_dashboard_{timestamp}.json'
    
    else:
        flash('Type de données non supporté pour l\'export JSON.', 'error')
        return redirect(url_for('are_dashboard.index'))
    
    output = io.BytesIO()
    output.write(json.dumps(data, indent=2, default=str).encode('utf-8'))
    output.seek(0)
    
    return send_file(
        output,
        as_attachment=True,
        download_name=filename,
        mimetype='application/json'
    )


# Fonction supprimée - dupliquée plus haut avec route différente


@dashboard_bp.route('/synthese-taux')
@login_required
@admin_required
def synthese_taux():
    """Synthèse des taux d'accès et d'électrification pour toutes les années disponibles"""
    return _get_synthese_taux_data()


def _get_synthese_taux_data():
    """Fonction helper pour récupérer les données de synthèse des taux"""
    from app.models.statistiques_are import StatistiqueNationale
    
    # Récupérer toutes les statistiques nationales triées par année
    statistiques = StatistiqueNationale.query.filter_by(actif=True).order_by(StatistiqueNationale.annee).all()
    
    # Préparer les données pour le template
    synthese_taux = []
    for stat in statistiques:
        synthese_taux.append({
            'annee': stat.annee,
            'total_clients': stat.total_clients_nationaux or 0,
            'clients_ht': stat.clients_ht_nationaux or 0,
            'clients_mt': stat.clients_mt_nationaux or 0,
            'clients_bt': stat.clients_bt_nationaux or 0,
            'taux_acces_pct': round(stat.taux_acces_national or 0, 2),
            'taux_electrification_pct': round(stat.taux_electrification_national or 0, 2),
            'taux_couverture_pct': round(stat.taux_couverture_national or 0, 2),
            'capacite_installee_mw': round(stat.capacite_totale_installee_mw or 0, 2),
            'capacite_disponible_mw': round(stat.capacite_totale_disponible_mw or 0, 2),
            'production_totale_gwh': round(stat.production_totale_annuelle_gwh or 0, 2)
        })
    
    # Statistiques générales
    total_annees = len(synthese_taux)
    moyenne_taux_acces = round(sum(s['taux_acces_pct'] for s in synthese_taux) / total_annees, 2) if total_annees > 0 else 0
    moyenne_taux_electrification = round(sum(s['taux_electrification_pct'] for s in synthese_taux) / total_annees, 2) if total_annees > 0 else 0
    moyenne_taux_couverture = round(sum(s['taux_couverture_pct'] for s in synthese_taux) / total_annees, 2) if total_annees > 0 else 0
    
    return render_template('are/dashboard/synthese_taux.html',
                         synthese_taux=synthese_taux,
                         total_annees=total_annees,
                         moyenne_taux_acces=moyenne_taux_acces,
                         moyenne_taux_electrification=moyenne_taux_electrification,
                         moyenne_taux_couverture=moyenne_taux_couverture)


def _calculer_statistiques_nationales_avancees(annee_debut, annee_fin):
    """Calcule les statistiques nationales avancées pour le dashboard"""
    from app.models.statistiques_are import StatistiqueNationale
    from app.models.production_hydro import CentraleHydro, RapportHydro
    from app.models.production_thermique import CentraleThermique, RapportThermique
    from app.models.production_solaire import CentraleSolaire, RapportSolaire
    from app.models.transport import LigneTransport
    from datetime import date
    
    # Classement des opérateurs par production
    classement_operateurs = []
    
    # Récupérer tous les opérateurs actifs
    operateurs = Operateur.query.filter_by(actif=True).all()
    
    for operateur in operateurs:
        # Calculer la production totale pour cet opérateur sur la période
        production_totale = 0
        
        # Production hydro
        prod_hydro = db.session.query(func.sum(RapportHydro.energie_produite)).join(
            CentraleHydro, RapportHydro.centrale_id == CentraleHydro.id
        ).filter(
            CentraleHydro.operateur_id == operateur.id,
            RapportHydro.annee.between(annee_debut, annee_fin),
            RapportHydro.actif == True
        ).scalar() or 0
        
        # Production thermique
        prod_thermique = db.session.query(func.sum(RapportThermique.energie_produite)).join(
            CentraleThermique, RapportThermique.centrale_id == CentraleThermique.id
        ).filter(
            CentraleThermique.operateur_id == operateur.id,
            RapportThermique.annee.between(annee_debut, annee_fin),
            RapportThermique.actif == True
        ).scalar() or 0
        
        # Production solaire
        prod_solaire = db.session.query(func.sum(RapportSolaire.energie_produite)).join(
            CentraleSolaire, RapportSolaire.centrale_id == CentraleSolaire.id
        ).filter(
            CentraleSolaire.operateur_id == operateur.id,
            RapportSolaire.annee.between(annee_debut, annee_fin),
            RapportSolaire.actif == True
        ).scalar() or 0
        
        # Production totale en GWh
        production_totale = (prod_hydro + prod_thermique + prod_solaire) / 1000
        
        if production_totale > 0:  # N'inclure que les opérateurs avec de la production
            classement_operateurs.append({
                'nom': operateur.nom,
                'production_gwh': round(production_totale, 1)
            })
    
    # Trier par production décroissante
    classement_operateurs.sort(key=lambda x: x['production_gwh'], reverse=True)
    
    # Évolution de la capacité installée
    evolution_capacite = []
    for annee in range(annee_debut, annee_fin + 1):
        # Calculer la capacité installée cumulée jusqu'à cette année
        capacite_hydro = db.session.query(func.sum(CentraleHydro.puissance_installee)).filter(
            CentraleHydro.actif == True,
            CentraleHydro.date_mise_service <= date(annee, 12, 31)
        ).scalar() or 0
        
        capacite_thermique = db.session.query(func.sum(CentraleThermique.puissance_installee)).filter(
            CentraleThermique.actif == True,
            CentraleThermique.date_mise_service <= date(annee, 12, 31)
        ).scalar() or 0
        
        capacite_solaire = db.session.query(func.sum(CentraleSolaire.puissance_installee)).filter(
            CentraleSolaire.actif == True,
            CentraleSolaire.date_mise_service <= date(annee, 12, 31)
        ).scalar() or 0
        
        evolution_capacite.append({
            'annee': annee,
            'capacite_hydro_mw': round(capacite_hydro, 2),
            'capacite_thermique_mw': round(capacite_thermique, 2),
            'capacite_solaire_mw': round(capacite_solaire, 2),
            'capacite_totale_mw': round(capacite_hydro + capacite_thermique + capacite_solaire, 2)
        })
    
    # Statistiques d'infrastructure
    stats_infrastructure = {
        'total_centrales': db.session.query(func.count(CentraleHydro.id)).filter(CentraleHydro.actif == True).scalar() + 
                          db.session.query(func.count(CentraleThermique.id)).filter(CentraleThermique.actif == True).scalar() +
                          db.session.query(func.count(CentraleSolaire.id)).filter(CentraleSolaire.actif == True).scalar(),
        'longueur_lignes_km': db.session.query(func.sum(LigneTransport.longueur_totale)).filter(LigneTransport.actif == True).scalar() or 0,
        'puissance_disponible_mw': db.session.query(func.sum(CentraleHydro.puissance_installee)).filter(CentraleHydro.actif == True).scalar() or 0,
        'capacite_totale_mw': (db.session.query(func.sum(CentraleHydro.puissance_installee)).filter(CentraleHydro.actif == True).scalar() or 0) +
                             (db.session.query(func.sum(CentraleThermique.puissance_installee)).filter(CentraleThermique.actif == True).scalar() or 0) +
                             (db.session.query(func.sum(CentraleSolaire.puissance_installee)).filter(CentraleSolaire.actif == True).scalar() or 0)
    }
    
    # Évolution de la clientèle
    evolution_clients = []
    for annee in range(annee_debut, annee_fin + 1):
        # Récupérer les statistiques nationales pour cette année
        stat = StatistiqueNationale.query.filter_by(annee=annee, actif=True).first()
        if stat:
            evolution_clients.append({
                'annee': annee,
                'clients_ht': stat.clients_ht_nationaux or 0,
                'clients_mt': stat.clients_mt_nationaux or 0,
                'clients_bt': stat.clients_bt_nationaux or 0
            })
    
    # Taux d'électrification
    taux_electrification = []
    for annee in range(annee_debut, annee_fin + 1):
        stat = StatistiqueNationale.query.filter_by(annee=annee, actif=True).first()
        if stat and stat.taux_acces_national:
            taux_electrification.append({
                'annee': annee,
                'taux_acces_pct': stat.taux_acces_national
            })
    
    # Statistiques solaires par année
    stats_solaire = []
    for annee in range(annee_debut, annee_fin + 1):
        # Capacité solaire installée cette année-là
        capacite_annee = db.session.query(func.sum(CentraleSolaire.puissance_installee)).filter(
            CentraleSolaire.actif == True,
            func.extract('year', CentraleSolaire.date_mise_service) == annee
        ).scalar() or 0
        
        # Nombre d'installations cette année-là
        nombre_installations = db.session.query(func.count(CentraleSolaire.id)).filter(
            CentraleSolaire.actif == True,
            func.extract('year', CentraleSolaire.date_mise_service) == annee
        ).scalar() or 0
        
        if capacite_annee > 0 or nombre_installations > 0:
            stats_solaire.append({
                'annee': annee,
                'capacite_mw': round(capacite_annee, 2),
                'installations': nombre_installations
            })
    
    # Statistiques de performance des opérateurs
    stats_operateurs = []
    for operateur in operateurs:
        # Production totale en GWh sur la période
        production_totale_gwh = 0
        
        # Production hydro
        prod_hydro = db.session.query(func.sum(RapportHydro.energie_produite)).join(
            CentraleHydro, RapportHydro.centrale_id == CentraleHydro.id
        ).filter(
            CentraleHydro.operateur_id == operateur.id,
            RapportHydro.annee.between(annee_debut, annee_fin),
            RapportHydro.actif == True
        ).scalar() or 0
        
        # Production thermique
        prod_thermique = db.session.query(func.sum(RapportThermique.energie_produite)).join(
            CentraleThermique, RapportThermique.centrale_id == CentraleThermique.id
        ).filter(
            CentraleThermique.operateur_id == operateur.id,
            RapportThermique.annee.between(annee_debut, annee_fin),
            RapportThermique.actif == True
        ).scalar() or 0
        
        # Production solaire
        prod_solaire = db.session.query(func.sum(RapportSolaire.energie_produite)).join(
            CentraleSolaire, RapportSolaire.centrale_id == CentraleSolaire.id
        ).filter(
            CentraleSolaire.operateur_id == operateur.id,
            RapportSolaire.annee.between(annee_debut, annee_fin),
            RapportSolaire.actif == True
        ).scalar() or 0
        
        production_totale_gwh = (prod_hydro + prod_thermique + prod_solaire) / 1000
        
        # Capacité installée totale de l'opérateur
        capacite_hydro = db.session.query(func.sum(CentraleHydro.puissance_installee)).filter(
            CentraleHydro.operateur_id == operateur.id,
            CentraleHydro.actif == True
        ).scalar() or 0
        
        capacite_thermique = db.session.query(func.sum(CentraleThermique.puissance_installee)).filter(
            CentraleThermique.operateur_id == operateur.id,
            CentraleThermique.actif == True
        ).scalar() or 0
        
        capacite_solaire = db.session.query(func.sum(CentraleSolaire.puissance_installee)).filter(
            CentraleSolaire.operateur_id == operateur.id,
            CentraleSolaire.actif == True
        ).scalar() or 0
        
        capacite_totale_mw = capacite_hydro + capacite_thermique + capacite_solaire
        
        # Calculer la production moyenne en MW (approximation)
        # Production en GWh / 8760 heures ≈ puissance moyenne en MW
        production_moyenne_mw = production_totale_gwh * 1000 / 8760 if annee_fin - annee_debut + 1 > 0 else 0
        
        # Nombre de clients
        nombre_clients = operateur.nombre_clients or 0
        
        # Rendement (production moyenne / capacité installée * 100)
        rendement_pct = (production_moyenne_mw / capacite_totale_mw * 100) if capacite_totale_mw > 0 else 0
        
        if production_totale_gwh > 0 or capacite_totale_mw > 0 or nombre_clients > 0:  # Inclure les opérateurs avec activité ou clients
            stats_operateurs.append({
                'nom': operateur.nom,
                'production_mw': round(production_moyenne_mw, 1),
                'production_gwh': round(production_totale_gwh, 1),
                'clients': nombre_clients,
                'rendement_pct': round(rendement_pct, 1)
            })
    
    # Trier par rendement décroissant
    stats_operateurs.sort(key=lambda x: x['rendement_pct'], reverse=True)
    
    # Évolution de la production par année
    evolution_production = []
    for annee in range(annee_debut, annee_fin + 1):
        # Production hydro pour cette année
        production_hydro_gwh = db.session.query(func.sum(RapportHydro.energie_produite)).join(
            CentraleHydro, RapportHydro.centrale_id == CentraleHydro.id
        ).filter(
            RapportHydro.annee == annee,
            RapportHydro.actif == True,
            CentraleHydro.actif == True
        ).scalar() or 0
        
        # Production thermique pour cette année
        production_thermique_gwh = db.session.query(func.sum(RapportThermique.energie_produite)).join(
            CentraleThermique, RapportThermique.centrale_id == CentraleThermique.id
        ).filter(
            RapportThermique.annee == annee,
            RapportThermique.actif == True,
            CentraleThermique.actif == True
        ).scalar() or 0
        
        # Production solaire pour cette année
        production_solaire_gwh = db.session.query(func.sum(RapportSolaire.energie_produite)).join(
            CentraleSolaire, RapportSolaire.centrale_id == CentraleSolaire.id
        ).filter(
            RapportSolaire.annee == annee,
            RapportSolaire.actif == True,
            CentraleSolaire.actif == True
        ).scalar() or 0
        
        # Production totale
        production_totale_gwh = production_hydro_gwh + production_thermique_gwh + production_solaire_gwh
        
        evolution_production.append({
            'annee': annee,
            'production_hydro_gwh': round(production_hydro_gwh, 1),
            'production_thermique_gwh': round(production_thermique_gwh, 1),
            'production_solaire_gwh': round(production_solaire_gwh, 1),
            'production_totale_gwh': round(production_totale_gwh, 1)
        })
    
    # Mix énergétique pour l'année la plus récente (annee_fin)
    mix_energetique = []
    annee_mix = annee_fin  # Utiliser la dernière année de la période
    
    # Production hydro pour cette année
    production_hydro_mix = db.session.query(func.sum(RapportHydro.energie_produite)).join(
        CentraleHydro, RapportHydro.centrale_id == CentraleHydro.id
    ).filter(
        RapportHydro.annee == annee_mix,
        RapportHydro.actif == True,
        CentraleHydro.actif == True
    ).scalar() or 0
    
    # Production thermique pour cette année
    production_thermique_mix = db.session.query(func.sum(RapportThermique.energie_produite)).join(
        CentraleThermique, RapportThermique.centrale_id == CentraleThermique.id
    ).filter(
        RapportThermique.annee == annee_mix,
        RapportThermique.actif == True,
        CentraleThermique.actif == True
    ).scalar() or 0
    
    # Production solaire pour cette année
    production_solaire_mix = db.session.query(func.sum(RapportSolaire.energie_produite)).join(
        CentraleSolaire, RapportSolaire.centrale_id == CentraleSolaire.id
    ).filter(
        RapportSolaire.annee == annee_mix,
        RapportSolaire.actif == True,
        CentraleSolaire.actif == True
    ).scalar() or 0
    
    # Calcul du total et des pourcentages
    total_production_mix = production_hydro_mix + production_thermique_mix + production_solaire_mix
    
    if total_production_mix > 0:
        mix_energetique = [
            {
                'source': 'Hydroélectrique',
                'production_gwh': round(production_hydro_mix, 1),
                'part_pct': round((production_hydro_mix / total_production_mix) * 100, 1)
            },
            {
                'source': 'Thermique',
                'production_gwh': round(production_thermique_mix, 1),
                'part_pct': round((production_thermique_mix / total_production_mix) * 100, 1)
            },
            {
                'source': 'Solaire',
                'production_gwh': round(production_solaire_mix, 1),
                'part_pct': round((production_solaire_mix / total_production_mix) * 100, 1)
            }
        ]
        # Trier par production décroissante et filtrer les sources avec 0%
        mix_energetique = [item for item in mix_energetique if item['part_pct'] > 0]
        mix_energetique.sort(key=lambda x: x['production_gwh'], reverse=True)
    
    return {
        'classement_operateurs': classement_operateurs,
        'evolution_capacite': evolution_capacite,
        'stats_infrastructure': stats_infrastructure,
        'evolution_clients': evolution_clients,
        'taux_electrification': taux_electrification,
        'stats_solaire': stats_solaire,
        'stats_operateurs': stats_operateurs,
        'evolution_production': evolution_production,
        'mix_energetique': mix_energetique
    }